#!/usr/bin/env python3
"""Render a .pptx to one PNG per slide (approximate layout) via python-pptx -> HTML -> Chrome headless.

Usage: python3 render_deck.py <deck.pptx> <out-dir>
Outputs slide-1.png ... slide-N.png in out-dir.
"""
import json
import os
import re
import subprocess
import sys
import tempfile
from pathlib import Path
from xml.etree import ElementTree as ET

from pptx import Presentation
from pptx.util import Emu

NS = {
    'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
    'c': 'http://schemas.openxmlformats.org/drawingml/2006/chart',
}

CHROME = '/Applications/Google Chrome.app/Contents/MacOS/Google Chrome'
PX_PER_EMU = 96 / 914400.0


def emu_px(v):
    return round(v * PX_PER_EMU, 1)


def shape_geometry(shape):
    try:
        return (
            emu_px(shape.left or 0), emu_px(shape.top or 0),
            emu_px(shape.width or 0), emu_px(shape.height or 0),
        )
    except Exception:
        return (0.0, 0.0, 100.0, 30.0)


def text_blocks(shape):
    blocks = []
    if not getattr(shape, 'has_text_frame', False):
        return blocks
    tf = shape.text_frame
    bottom = 0.0
    for para in tf.paragraphs:
        text = ''.join(run.text for run in para.runs)
        if not text.strip():
            continue
        sizes = [run.font.size.pt if run.font.size else 18 for run in para.runs]
        size = max(sizes) if sizes else 18
        blocks.append({'text': text, 'size': max(9, min(size, 48))})
    return blocks


def chart_data(shape):
    """Extract categories + series values from a native chart in the shape."""
    try:
        chart = shape.chart
    except Exception:
        return None
    try:
        cats = list(chart.plots[0].categories)
        ser = chart.plots[0].series[0]
        vals = list(ser.values)
        title = ''
        try:
            title = chart.chart_title.text_frame.text or ''
        except Exception:
            pass
        return {'title': title, 'cats': [str(c) for c in cats], 'vals': [float(v) for v in vals]}
    except Exception:
        return None


def svg_chart(data, w, h):
    n = max(1, len(data['vals']))
    title_h = 28
    plot_h = h - title_h - 26
    col_w = w / n
    max_v = max(max(data['vals']), 1e-6)
    parts = [f'<svg width="{w}" height="{h}" xmlns="http://www.w3.org/2000/svg">',
             f'<rect width="{w}" height="{h}" fill="#f7f9fc"/>']
    if data['title']:
        parts.append(f'<text x="6" y="18" font-size="13" fill="#222" font-family="sans-serif">{html_escape(data["title"])}</text>')
    parts.append(f'<line x1="6" y1="{h-10}" x2="{w-4}" y2="{h-10}" stroke="#888" stroke-width="1"/>')
    for i, v in enumerate(data['vals']):
        x0 = i * col_w
        bh = max(2.0, (v / max_v) * plot_h)
        y0 = h - 10 - bh
        parts.append(f'<rect x="{x0+4}" y="{y0}" width="{col_w-8}" height="{bh}" fill="#3f7fd1"/>')
        parts.append(f'<text x="{x0+4}" y="{y0-3}" font-size="10" fill="#333" font-family="sans-serif">{v:.3g}</text>')
        cat = str(data['cats'][i]) if i < len(data['cats']) else ''
        if cat:
            parts.append(f'<text x="{x0+2}" y="{h-2}" font-size="9" fill="#555" font-family="sans-serif">{html_escape(cat[:16])}</text>')
    parts.append('</svg>')
    return '\n'.join(parts)


def page_html(slide, page_w, page_h):
    shapes = []
    charts = []
    for shape in slide.shapes:
        x, y, w, h = shape_geometry(shape)
        cd = chart_data(shape)
        if cd is not None:
            charts.append(
                f'<div style="position:absolute;left:{x}px;top:{y}px;width:{w}px;height:{h}px;">{svg_chart(cd, int(w), int(h))}</div>')
            continue
        blocks = text_blocks(shape)
        if not blocks:
            continue
        top = y
        for b in blocks[:10]:
            font = b['size']
            line_h = font * 1.25
            shapes.append(
                f'<div style="position:absolute;left:{x}px;top:{top}px;width:{w}px;'
                f'height:{line_h}px;font-size:{font}px;line-height:{line_h}px;color:#222;'
                f'font-family:sans-serif;white-space:pre-wrap;word-break:break-word;">{html_escape(b["text"])}</div>')
            top += line_h
    head = '<!doctype html><html><head><meta charset="utf-8"><style>body{margin:0;background:#fff;}</style></head>'
    body = f'<body style="width:{page_w}px;height:{page_h}px;position:relative;overflow:hidden;">'
    return head + body + ''.join(shapes) + ''.join(charts) + '</body></html>'


def pdf_html(pages, page_w, page_h):
    parts = [
        '<!doctype html><html><head><meta charset="utf-8"><style>'
        '@page { size: %dpx %dpx; margin: 0; }'
        'body { margin: 0; background: #fff; }'
        '.page { page-break-after: always; width: %dpx; height: %dpx; position: relative; overflow: hidden; }'
        '.page:last-child { page-break-after: auto; }'
        '</style></head><body>' % (page_w, page_h, page_w, page_h)
    ]
    for html in pages:
        body = html.split('<body', 1)[1]
        body = body.rsplit('</body>', 1)[0]
        parts.append(f'<div class="page">{body}')
        parts.append('</div>')
    parts.append('</body></html>')
    return ''.join(parts)


def html_escape(text):
    return (
        text.replace('&', '&amp;').replace('<', '&lt;').replace('>', '&gt;')
        .replace('"', '&quot;')
    )


def main():
    if len(sys.argv) < 3:
        print('usage: render_deck.py <deck.pptx> <out-dir>', file=sys.stderr)
        return 2
    deck = Path(sys.argv[1]).resolve()
    out = Path(sys.argv[2]).resolve()
    out.mkdir(parents=True, exist_ok=True)
    prs = Presentation(str(deck))
    page_w = emu_px(prs.slide_width)
    page_h = emu_px(prs.slide_height)
    with tempfile.TemporaryDirectory() as tmp:
        all_pages = []
        for idx, slide in enumerate(prs.slides, 1):
            html = page_html(slide, int(page_w), int(page_h))
            all_pages.append(html)
            html_path = Path(tmp) / f'slide-{idx}.html'
            html_path.write_text(html, encoding='utf-8')
            png = out / f'slide-{idx}.png'
            subprocess.run([
                CHROME, '--headless=new', '--disable-gpu', '--hide-scrollbars',
                f'--window-size={int(page_w)},{int(page_h)}',
                f'--screenshot={png}',
                html_path.as_uri(),
            ], capture_output=True, timeout=60, check=False)
            print(f'slide-{idx}.png {png.stat().st_size if png.exists() else "FAILED"}')
        # Also produce a combined printable PDF at deck.pdf (approximate visuals,
        # same page geometry) so a visual-capable reader can open one document.
        combined = Path(tmp) / 'deck.html'
        combined.write_text(pdf_html(all_pages, int(page_w), int(page_h)), encoding='utf-8')
        pdf = out / 'deck.pdf'
        subprocess.run([
            CHROME, '--headless=new', '--disable-gpu', '--no-pdf-header-footer',
            f'--print-to-pdf={pdf}',
            combined.as_uri(),
        ], capture_output=True, timeout=60, check=False)
        print(f'deck.pdf {pdf.stat().st_size if pdf.exists() else "FAILED"}')
    return 0


if __name__ == '__main__':
    sys.exit(main())
